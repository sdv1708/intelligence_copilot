"""Document parsing for PDF, DOCX, PPTX, and TXT files."""

from typing import Optional
import io
from core.utils import log_message


def parse_pdf(file_content: bytes) -> str:
    """Extract text from PDF file."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_content))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        log_message("INFO", f"Parsed PDF: {len(text)} characters")
        return text.strip()
    except Exception as e:
        log_message("ERROR", f"Failed to parse PDF: {str(e)}")
        return ""


def parse_docx(file_content: bytes) -> str:
    """Extract text from DOCX file."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_content))
        text = "\n".join([para.text for para in doc.paragraphs])
        log_message("INFO", f"Parsed DOCX: {len(text)} characters")
        return text.strip()
    except Exception as e:
        log_message("ERROR", f"Failed to parse DOCX: {str(e)}")
        return ""


def parse_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        log_message("INFO", f"Parsed PPTX: {len(text)} characters")
        return text.strip()
    except Exception as e:
        log_message("ERROR", f"Failed to parse PPTX: {str(e)}")
        return ""


def parse_txt(file_content: bytes) -> str:
    """Extract text from TXT file."""
    try:
        text = file_content.decode("utf-8", errors="ignore").strip()
        log_message("INFO", f"Parsed TXT: {len(text)} characters")
        return text
    except Exception as e:
        log_message("ERROR", f"Failed to parse TXT: {str(e)}")
        return ""


def parse_file(file_content: bytes, filename: str) -> tuple[str, str]:
    """
    Parse a file based on its extension.
    Returns (text, media_type).
    """
    filename_lower = filename.lower()
    
    if filename_lower.endswith(".pdf"):
        return parse_pdf(file_content), "pdf"
    elif filename_lower.endswith(".docx"):
        return parse_docx(file_content), "docx"
    elif filename_lower.endswith(".pptx"):
        return parse_pptx(file_content), "pptx"
    elif filename_lower.endswith(".txt"):
        return parse_txt(file_content), "txt"
    else:
        log_message("WARNING", f"Unsupported file type: {filename}")
        return "", "unknown"


def parse_pasted_text(text: str) -> tuple[str, str]:
    """Parse pasted text from textarea."""
    text = text.strip()
    if not text:
        log_message("WARNING", "Empty pasted text")
        return "", "pasted"
    log_message("INFO", f"Parsed pasted text: {len(text)} characters")
    return text, "pasted"

